"""
Export module for generating reports in multiple formats (PDF, DOCX, PPTX).
"""

from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime
from io import BytesIO

import pandas as pd

from src.config import PRODUCTS_DIR, REPORTS_DIR
from src.utils import setup_logging


class ReportExporter:
    """Exports reports to PDF, DOCX, and PPTX formats."""

    def __init__(self):
        self.logger = setup_logging("exporter")
        self.reports_dir = REPORTS_DIR
        self.reports_dir.mkdir(parents=True, exist_ok=True)

    def export_all(
        self,
        analysis_results: Dict[str, Any],
        ai_insights: Dict[str, Any],
        df: pd.DataFrame,
        visualisation_paths: Dict[str, Path],
        version: str = "v1"
    ) -> Dict[str, Path]:
        """Export report to all formats."""
        outputs = {}

        outputs["pdf"] = self.export_to_pdf(
            analysis_results, ai_insights, df, visualisation_paths, version
        )
        outputs["docx"] = self.export_to_docx(
            analysis_results, ai_insights, df, visualisation_paths, version
        )
        outputs["pptx"] = self.export_to_pptx(
            analysis_results, ai_insights, df, visualisation_paths, version
        )

        self.logger.info(f"Exported reports to {len(outputs)} formats")
        return outputs

    def export_to_pdf(
        self,
        analysis_results: Dict[str, Any],
        ai_insights: Dict[str, Any],
        df: pd.DataFrame,
        visualisation_paths: Dict[str, Path],
        version: str = "v1"
    ) -> Path:
        """Export report to PDF format using reportlab."""
        self.logger.info("Exporting to PDF")

        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            Image, PageBreak, ListFlowable, ListItem
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        output_path = self.reports_dir / f"churn_analysis_report_{version}.pdf"

        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='Title2',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            name='Subtitle',
            parent=styles['Normal'],
            fontSize=14,
            spaceAfter=12,
            alignment=TA_CENTER,
            textColor=colors.grey
        ))

        story = []

        # Title page
        story.append(Spacer(1, 2 * inch))
        story.append(Paragraph("Veritly Market Intelligence Report", styles['Title2']))
        story.append(Paragraph("Telecom Customer Churn Analysis", styles['Subtitle']))
        story.append(Spacer(1, 0.5 * inch))
        story.append(Paragraph(
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            styles['Subtitle']
        ))
        story.append(PageBreak())

        # Executive Summary
        story.append(Paragraph("Executive Summary", styles['Heading1']))
        story.append(Spacer(1, 12))

        churn_rate = analysis_results.get("overall_churn_rate", 0)
        total = analysis_results.get("total_count", 0)
        churned = analysis_results.get("churned_count", 0)

        summary_text = f"""
        Analysis of <b>{total:,}</b> telecom customers reveals a churn rate of <b>{churn_rate}%</b>,
        representing <b>{churned:,}</b> customers lost. Key findings indicate that international plan
        customers and those with high service call volumes are at elevated risk of churning.
        """
        story.append(Paragraph(summary_text, styles['Normal']))
        story.append(Spacer(1, 24))

        # Key Metrics Table
        story.append(Paragraph("Key Metrics", styles['Heading2']))
        story.append(Spacer(1, 12))

        validation = analysis_results.get("model_validation", {})
        metrics_data = [
            ["Metric", "Value"],
            ["Total Customers", f"{total:,}"],
            ["Churned Customers", f"{churned:,}"],
            ["Churn Rate", f"{churn_rate}%"],
            ["Model Accuracy", f"{validation.get('accuracy', 'N/A')}%"],
            ["High Risk Customers", f"{validation.get('high_risk_count', 0):,}"]
        ]

        metrics_table = Table(metrics_data, colWidths=[3 * inch, 2 * inch])
        metrics_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498db')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#ecf0f1')),
            ('GRID', (0, 0), (-1, -1), 1, colors.white),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
        ]))
        story.append(metrics_table)
        story.append(Spacer(1, 24))

        # Visualisations
        story.append(PageBreak())
        story.append(Paragraph("Analysis Visualisations", styles['Heading1']))
        story.append(Spacer(1, 12))

        for chart_name, chart_path in visualisation_paths.items():
            if chart_path and chart_path.exists():
                story.append(Paragraph(
                    chart_name.replace('_', ' ').title(),
                    styles['Heading3']
                ))
                story.append(Spacer(1, 6))
                img = Image(str(chart_path), width=5.5 * inch, height=3.5 * inch)
                story.append(img)
                story.append(Spacer(1, 24))

        # Risk Tiers
        story.append(PageBreak())
        story.append(Paragraph("Risk Tier Distribution", styles['Heading1']))
        story.append(Spacer(1, 12))

        tier_dist = validation.get("risk_tier_distribution", {})
        risk_data = [
            ["Risk Tier", "Count", "Probability Range"],
            ["Low", f"{tier_dist.get('low', 0):,}", "0% - 15%"],
            ["Medium", f"{tier_dist.get('medium', 0):,}", "15% - 30%"],
            ["High", f"{tier_dist.get('high', 0):,}", "30% - 50%"],
            ["Critical", f"{tier_dist.get('critical', 0):,}", "50% - 100%"]
        ]

        risk_table = Table(risk_data, colWidths=[2 * inch, 2 * inch, 2 * inch])
        risk_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c3e50')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (0, 1), colors.HexColor('#27ae60')),
            ('BACKGROUND', (0, 2), (0, 2), colors.HexColor('#f1c40f')),
            ('BACKGROUND', (0, 3), (0, 3), colors.HexColor('#e67e22')),
            ('BACKGROUND', (0, 4), (0, 4), colors.HexColor('#e74c3c')),
            ('TEXTCOLOR', (0, 1), (0, -1), colors.white),
            ('GRID', (0, 0), (-1, -1), 1, colors.white),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
        ]))
        story.append(risk_table)

        # Footer
        story.append(Spacer(1, inch))
        story.append(Paragraph(
            "Generated by Veritly AI Market Intelligence Platform",
            styles['Subtitle']
        ))

        doc.build(story)
        self.logger.info(f"PDF exported to: {output_path}")
        return output_path

    def export_to_docx(
        self,
        analysis_results: Dict[str, Any],
        ai_insights: Dict[str, Any],
        df: pd.DataFrame,
        visualisation_paths: Dict[str, Path],
        version: str = "v1"
    ) -> Path:
        """Export report to Word document format."""
        self.logger.info("Exporting to DOCX")

        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT

        doc = Document()

        # Title
        title = doc.add_heading('Veritly Market Intelligence Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        subtitle = doc.add_paragraph('Telecom Customer Churn Analysis')
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        doc.add_paragraph()

        # Executive Summary
        doc.add_heading('Executive Summary', level=1)

        churn_rate = analysis_results.get("overall_churn_rate", 0)
        total = analysis_results.get("total_count", 0)
        churned = analysis_results.get("churned_count", 0)

        summary = doc.add_paragraph()
        summary.add_run(f"Analysis of {total:,} telecom customers reveals a churn rate of ")
        summary.add_run(f"{churn_rate}%").bold = True
        summary.add_run(f", representing {churned:,} customers lost.")

        doc.add_paragraph(
            "Key findings indicate that international plan customers and those with "
            "high service call volumes are at elevated risk of churning."
        )

        # Key Metrics
        doc.add_heading('Key Metrics', level=1)

        validation = analysis_results.get("model_validation", {})

        table = doc.add_table(rows=6, cols=2)
        table.style = 'Table Grid'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        metrics = [
            ("Metric", "Value"),
            ("Total Customers", f"{total:,}"),
            ("Churned Customers", f"{churned:,}"),
            ("Churn Rate", f"{churn_rate}%"),
            ("Model Accuracy", f"{validation.get('accuracy', 'N/A')}%"),
            ("High Risk Customers", f"{validation.get('high_risk_count', 0):,}")
        ]

        for i, (metric, value) in enumerate(metrics):
            row = table.rows[i]
            row.cells[0].text = metric
            row.cells[1].text = value
            if i == 0:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True

        doc.add_paragraph()

        # Segment Analysis
        doc.add_heading('Segment Analysis', level=1)

        segments = analysis_results.get("segment_analysis", {})

        # International Plan
        doc.add_heading('By International Plan', level=2)
        intl_data = segments.get("international_plan", {})
        if intl_data:
            table = doc.add_table(rows=len(intl_data) + 1, cols=4)
            table.style = 'Table Grid'
            hdr = table.rows[0].cells
            hdr[0].text = "Segment"
            hdr[1].text = "Customers"
            hdr[2].text = "Churned"
            hdr[3].text = "Churn Rate"
            for i, (segment, data) in enumerate(intl_data.items(), 1):
                row = table.rows[i].cells
                row[0].text = segment
                row[1].text = f"{data['total']:,}"
                row[2].text = f"{data['churned']:,}"
                row[3].text = f"{data['churn_rate']}%"

        doc.add_paragraph()

        # Visualisations
        doc.add_heading('Analysis Visualisations', level=1)

        for chart_name, chart_path in visualisation_paths.items():
            if chart_path and chart_path.exists():
                doc.add_heading(chart_name.replace('_', ' ').title(), level=2)
                doc.add_picture(str(chart_path), width=Inches(5.5))
                doc.add_paragraph()

        # Risk Distribution
        doc.add_heading('Risk Tier Distribution', level=1)

        tier_dist = validation.get("risk_tier_distribution", {})
        table = doc.add_table(rows=5, cols=3)
        table.style = 'Table Grid'

        risk_rows = [
            ("Risk Tier", "Count", "Probability Range"),
            ("Low", f"{tier_dist.get('low', 0):,}", "0% - 15%"),
            ("Medium", f"{tier_dist.get('medium', 0):,}", "15% - 30%"),
            ("High", f"{tier_dist.get('high', 0):,}", "30% - 50%"),
            ("Critical", f"{tier_dist.get('critical', 0):,}", "50% - 100%")
        ]

        for i, (tier, count, range_) in enumerate(risk_rows):
            row = table.rows[i].cells
            row[0].text = tier
            row[1].text = count
            row[2].text = range_

        # Footer
        doc.add_paragraph()
        footer = doc.add_paragraph("Generated by Veritly AI Market Intelligence Platform")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

        output_path = self.reports_dir / f"churn_analysis_report_{version}.docx"
        doc.save(str(output_path))
        self.logger.info(f"DOCX exported to: {output_path}")
        return output_path

    def export_to_pptx(
        self,
        analysis_results: Dict[str, Any],
        ai_insights: Dict[str, Any],
        df: pd.DataFrame,
        visualisation_paths: Dict[str, Path],
        version: str = "v1"
    ) -> Path:
        """Export report to PowerPoint presentation."""
        self.logger.info("Exporting to PPTX")

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title Slide
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Veritly Market Intelligence Report"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Telecom Customer Churn Analysis"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(127, 140, 141)
        p.alignment = PP_ALIGN.CENTER

        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime('%Y-%m-%d')
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(149, 165, 166)
        p.alignment = PP_ALIGN.CENTER

        # Executive Summary Slide
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)

        churn_rate = analysis_results.get("overall_churn_rate", 0)
        total = analysis_results.get("total_count", 0)
        churned = analysis_results.get("churned_count", 0)
        validation = analysis_results.get("model_validation", {})

        # Key metrics boxes
        metrics = [
            ("Total Customers", f"{total:,}", RGBColor(52, 152, 219)),
            ("Churn Rate", f"{churn_rate}%", RGBColor(231, 76, 60)),
            ("Churned", f"{churned:,}", RGBColor(230, 126, 34)),
            ("Model Accuracy", f"{validation.get('accuracy', 'N/A')}%", RGBColor(46, 204, 113))
        ]

        for i, (label, value, color) in enumerate(metrics):
            left = Inches(0.5 + i * 3.2)
            top = Inches(1.5)

            shape = slide.shapes.add_shape(1, left, top, Inches(2.9), Inches(1.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            # Value
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Label
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Key findings
        findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(3))
        tf = findings_box.text_frame
        tf.word_wrap = True

        findings = [
            "International plan customers show significantly elevated churn rates",
            "Customer service calls are a key predictor of churn risk",
            "Voicemail plan subscribers demonstrate lower churn rates",
            f"Model identifies {validation.get('high_risk_count', 0):,} high-risk customers"
        ]

        p = tf.paragraphs[0]
        p.text = "Key Findings:"
        p.font.size = Pt(20)
        p.font.bold = True

        for finding in findings:
            p = tf.add_paragraph()
            p.text = f"  {finding}"
            p.font.size = Pt(16)
            p.level = 0

        # Visualisation Slides
        for chart_name, chart_path in visualisation_paths.items():
            if chart_path and chart_path.exists():
                slide = prs.slides.add_slide(slide_layout)

                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = chart_name.replace('_', ' ').title()
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)

                # Chart image
                slide.shapes.add_picture(
                    str(chart_path),
                    Inches(1.5), Inches(1.3),
                    width=Inches(10), height=Inches(5.5)
                )

        # Risk Distribution Slide
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Risk Tier Distribution"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)

        tier_dist = validation.get("risk_tier_distribution", {})
        tiers = [
            ("Low Risk", tier_dist.get('low', 0), RGBColor(39, 174, 96)),
            ("Medium Risk", tier_dist.get('medium', 0), RGBColor(241, 196, 15)),
            ("High Risk", tier_dist.get('high', 0), RGBColor(230, 126, 34)),
            ("Critical", tier_dist.get('critical', 0), RGBColor(231, 76, 60))
        ]

        for i, (label, count, color) in enumerate(tiers):
            left = Inches(0.8 + i * 3.1)
            top = Inches(2)

            shape = slide.shapes.add_shape(1, left, top, Inches(2.8), Inches(2.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{count:,}"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Closing Slide
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Veritly AI"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Market Intelligence Platform"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(127, 140, 141)
        p.alignment = PP_ALIGN.CENTER

        output_path = self.reports_dir / f"churn_analysis_report_{version}.pptx"
        prs.save(str(output_path))
        self.logger.info(f"PPTX exported to: {output_path}")
        return output_path
