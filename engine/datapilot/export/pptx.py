"""
PPTX export — generate PowerPoint presentations from analysis results.

Uses python-pptx. All text is configurable via parameters.
"""

import re
from datetime import datetime
from pathlib import Path
from typing import Any


def _split_sentences(text: str) -> list[str]:
    """Split text into sentences on sentence-ending . or ; boundaries."""
    parts = re.split(r'(?<!\d)\.(?!\d)\s*|;\s*', text.strip())
    return [s.strip() for s in parts if s.strip()]


def _truncate_words(text: str, max_words: int) -> str:
    """Truncate text to max_words."""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "..."


def export_to_pptx(
    analysis_results: dict[str, Any],
    output_path: str,
    title: str = "Data Analysis Report",
    subtitle: str = "Comprehensive Analysis",
    brand_name: str = "DataPilot",
    visualisation_paths: dict[str, Path] | None = None,
    metrics: list[dict[str, str]] | None = None,
) -> str:
    """
    Export analysis results to PowerPoint presentation.

    Args:
        analysis_results: Dict with analysis output (summary, sections, key_points, metrics).
        output_path: Where to save the PPTX.
        title: Report title.
        subtitle: Report subtitle.
        brand_name: Brand name for closing slide.
        visualisation_paths: Dict of {name: path} for chart images.
        metrics: List of {"label": ..., "value": ...} dicts for key metrics.

    Returns:
        The output file path.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    p = Path(output_path)
    p.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank

    # --- Title Slide ---
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    pg = tf.paragraphs[0]
    pg.text = title
    pg.font.size = Pt(44)
    pg.font.bold = True
    pg.font.color.rgb = RGBColor(44, 62, 80)
    pg.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    pg = tf.paragraphs[0]
    pg.text = subtitle
    pg.font.size = Pt(28)
    pg.font.color.rgb = RGBColor(127, 140, 141)
    pg.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    pg = tf.paragraphs[0]
    pg.text = datetime.now().strftime('%Y-%m-%d')
    pg.font.size = Pt(18)
    pg.font.color.rgb = RGBColor(149, 165, 166)
    pg.alignment = PP_ALIGN.CENTER

    # --- Executive Summary Slide (bullet points) ---
    summary_text = analysis_results.get("summary", "")
    if summary_text:
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        pg = tf.paragraphs[0]
        pg.text = "Executive Summary"
        pg.font.size = Pt(36)
        pg.font.bold = True
        pg.font.color.rgb = RGBColor(44, 62, 80)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5))
        tf = body_box.text_frame
        tf.word_wrap = True

        # Convert summary to max 5 bullet points
        sentences = _split_sentences(summary_text)
        bullets = [_truncate_words(s, 20) for s in sentences[:5]]
        if not bullets:
            bullets = [_truncate_words(summary_text, 20)]

        for j, bullet in enumerate(bullets):
            if j == 0:
                pg = tf.paragraphs[0]
            else:
                pg = tf.add_paragraph()
            pg.text = f"\u2022 {bullet}"
            pg.font.size = Pt(18)
            pg.font.color.rgb = RGBColor(52, 73, 94)

    # --- Metrics Slide (up to 8 metrics, 2 rows x 4 cols) ---
    report_metrics = metrics or analysis_results.get("metrics")
    if report_metrics:
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        pg = tf.paragraphs[0]
        pg.text = "Key Metrics"
        pg.font.size = Pt(36)
        pg.font.bold = True
        pg.font.color.rgb = RGBColor(44, 62, 80)

        colors_list = [
            RGBColor(52, 152, 219),   # blue
            RGBColor(231, 76, 60),    # red
            RGBColor(230, 126, 34),   # orange
            RGBColor(46, 204, 113),   # green
            RGBColor(155, 89, 182),   # purple
            RGBColor(26, 188, 156),   # teal
            RGBColor(241, 196, 15),   # yellow
            RGBColor(44, 62, 80),     # dark blue
        ]

        display_metrics = report_metrics[:8]
        items_per_row = min(4, len(display_metrics))
        box_width = min(2.9, (12.333 - 0.5) / items_per_row - 0.3)

        for i, m in enumerate(display_metrics):
            col = i % 4
            row = i // 4
            left = Inches(0.5 + col * (box_width + 0.3))
            top = Inches(1.5 + row * 2.2)
            color = colors_list[i % len(colors_list)]

            shape = slide.shapes.add_shape(1, left, top, Inches(box_width), Inches(1.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            tf = shape.text_frame
            tf.clear()
            pg = tf.paragraphs[0]
            pg.text = str(m.get("value", ""))
            pg.font.size = Pt(32 if row > 0 or len(display_metrics) > 4 else 40)
            pg.font.bold = True
            pg.font.color.rgb = RGBColor(255, 255, 255)
            pg.alignment = PP_ALIGN.CENTER

            pg = tf.add_paragraph()
            pg.text = m.get("label", "")
            pg.font.size = Pt(14)
            pg.font.color.rgb = RGBColor(255, 255, 255)
            pg.alignment = PP_ALIGN.CENTER

    # --- Analysis Section Slides (all bullet format, inline charts) ---
    sections = analysis_results.get("sections", [])
    section_skills_with_charts = set()
    for section in sections:
        slide = prs.slides.add_slide(slide_layout)

        heading = section.get("heading", "Analysis")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        pg = tf.paragraphs[0]
        pg.text = heading
        pg.font.size = Pt(32)
        pg.font.bold = True
        pg.font.color.rgb = RGBColor(44, 62, 80)

        # Collect bullets: key_points first, then extract from narrative
        bullets: list[str] = []
        key_points = section.get("key_points", [])
        for kp in key_points:
            bullets.append(_truncate_words(kp, 15))

        narrative = section.get("narrative", "")
        if len(bullets) < 6 and narrative:
            sentences = _split_sentences(narrative)
            for s in sentences:
                if len(bullets) >= 6:
                    break
                truncated = _truncate_words(s, 15)
                if truncated not in bullets:
                    bullets.append(truncated)

        bullets = bullets[:6]

        # Check if chart exists for this section
        skill = section.get("skill", "")
        has_chart = False
        chart_cp = None
        if visualisation_paths and skill in visualisation_paths:
            chart_cp = Path(visualisation_paths[skill]) if not isinstance(visualisation_paths[skill], Path) else visualisation_paths[skill]
            if chart_cp.exists():
                has_chart = True
                section_skills_with_charts.add(skill)

        # Layout: text area width depends on chart presence
        text_width = Inches(5.5) if has_chart else Inches(11.733)
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), text_width, Inches(5.5))
        tf = body_box.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            if j == 0:
                pg = tf.paragraphs[0]
            else:
                pg = tf.add_paragraph()
            pg.text = f"\u2022 {bullet}"
            pg.font.size = Pt(16)
            pg.font.color.rgb = RGBColor(52, 73, 94)

        # Add inline chart on right half
        if has_chart and chart_cp:
            slide.shapes.add_picture(
                str(chart_cp),
                Inches(6.8), Inches(1.3),
                width=Inches(6), height=Inches(5.5),
            )

    # --- Catch-all for unmatched chart_paths ---
    if visualisation_paths:
        unmatched = {k: v for k, v in visualisation_paths.items() if k not in section_skills_with_charts}
        for chart_name, chart_path in unmatched.items():
            cp = Path(chart_path) if not isinstance(chart_path, Path) else chart_path
            if cp.exists():
                slide = prs.slides.add_slide(slide_layout)

                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
                tf = title_box.text_frame
                pg = tf.paragraphs[0]
                pg.text = chart_name.replace('_', ' ').title()
                pg.font.size = Pt(32)
                pg.font.bold = True
                pg.font.color.rgb = RGBColor(44, 62, 80)

                slide.shapes.add_picture(
                    str(cp),
                    Inches(1.5), Inches(1.3),
                    width=Inches(10), height=Inches(5.5),
                )

    # --- Closing Slide with Key Recommendations ---
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    pg = tf.paragraphs[0]
    pg.text = brand_name
    pg.font.size = Pt(44)
    pg.font.bold = True
    pg.font.color.rgb = RGBColor(44, 62, 80)
    pg.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.6))
    tf = subtitle_box.text_frame
    pg = tf.paragraphs[0]
    pg.text = "AI-Powered Data Analysis"
    pg.font.size = Pt(24)
    pg.font.color.rgb = RGBColor(127, 140, 141)
    pg.alignment = PP_ALIGN.CENTER

    # Key Recommendations from top key_points
    all_key_points = analysis_results.get("key_points", [])
    if all_key_points:
        rec_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(2.5))
        tf = rec_box.text_frame
        tf.word_wrap = True

        rec_title = tf.paragraphs[0]
        rec_title.text = "Key Recommendations"
        rec_title.font.size = Pt(20)
        rec_title.font.bold = True
        rec_title.font.color.rgb = RGBColor(44, 62, 80)
        rec_title.alignment = PP_ALIGN.CENTER

        for kp in all_key_points[:3]:
            pg = tf.add_paragraph()
            pg.text = f"\u2022 {_truncate_words(kp, 15)}"
            pg.font.size = Pt(16)
            pg.font.color.rgb = RGBColor(52, 73, 94)
            pg.alignment = PP_ALIGN.LEFT

    prs.save(str(p))
    return str(p)
